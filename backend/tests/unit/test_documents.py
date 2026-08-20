"""Unit tests for document ingestion modules."""

from __future__ import annotations

import asyncio

import pytest
from fastapi import HTTPException

from app.documents.entity_extraction import extract_entities
from app.documents.embedding import is_broken_vector
from app.documents.validation import read_upload, sanitize_filename, validate_upload
from app.documents.chunking.structural_chunker import StructuralChunker
from app.documents.text_extraction import ExtractedText, extract_text
from app.documents.metadata_extraction import extract_metadata
class TestValidation:
    def test_valid_txt_file(self):
        result = validate_upload("doc.txt", 1024)
        assert result.valid is True
        assert result.error is None

    def test_valid_pdf_file(self):
        result = validate_upload("doc.pdf", 1024 * 1024)
        assert result.valid is True

    def test_invalid_extension(self):
        result = validate_upload("doc.exe", 100)
        assert result.valid is False
        assert "not allowed" in result.error

    def test_oversized_file(self):
        result = validate_upload("doc.txt", 25 * 1024 * 1024)
        assert result.valid is False
        assert "too big" in result.error
        assert "20 MB" in result.error
        assert "smaller than" in result.error

    def test_case_insensitive_extension(self):
        result = validate_upload("DOC.PDF", 1024)
        assert result.valid is True


class TestSanitizeFilename:
    def test_plain_name_passes(self):
        assert sanitize_filename("loan-letter.pdf") == "loan-letter.pdf"
        assert sanitize_filename("my doc (2).PDF") == "my doc (2).PDF"

    def test_rejects_empty_and_dots(self):
        assert sanitize_filename("") is None
        assert sanitize_filename(".") is None
        assert sanitize_filename("..") is None

    def test_rejects_forward_slash_traversal(self):
        assert sanitize_filename("../../../evil.pdf") is None
        assert sanitize_filename("a/b.pdf") is None

    def test_rejects_backslash_traversal(self):
        assert sanitize_filename("..\\..\\evil.pdf") is None
        assert sanitize_filename("sub\\b.pdf") is None

    def test_traversal_name_fails_validation(self):
        result = validate_upload("../../../evil.pdf", 100)
        assert result.valid is False
        assert "unsafe" in result.error

    def test_traversal_name_fails_validation_windows(self):
        result = validate_upload("..\\..\\evil.pdf", 100)
        assert result.valid is False


class TestReadUpload:
    class _FakeUpload:
        def __init__(self, chunks: list[bytes]):
            self._chunks = iter(chunks)

        async def read(self, size: int):
            try:
                return next(self._chunks)
            except StopIteration:
                return b""

    def test_reads_full_body(self):
        upload = self._FakeUpload([b"a" * 8192, b"b" * 8192, b"c"])
        total, content = asyncio.run(read_upload(upload, max_bytes=10_000_000))
        assert total == 8192 * 2 + 1
        assert content == b"a" * 8192 + b"b" * 8192 + b"c"

    def test_aborts_413_when_body_exceeds_cap(self):
        upload = self._FakeUpload([b"x" * 8192] * 5)
        with pytest.raises(HTTPException) as exc:
            asyncio.run(read_upload(upload, max_bytes=20_000))
        assert exc.value.status_code == 413

    def test_cap_defaults_to_settings(self):
        upload = self._FakeUpload([b"x" * 8192] * 5)
        # 5 * 8192 = 40 KB > default 20 MB? No — below it, so it must succeed.
        total, _ = asyncio.run(read_upload(upload))
        assert total == 8192 * 5


class TestBrokenVector:
    def test_missing_vector_is_broken(self):
        assert is_broken_vector(None) is True
        assert is_broken_vector([]) is True

    def test_zero_vector_is_broken(self):
        assert is_broken_vector([0.0] * 384) is True

    def test_nan_vector_is_broken(self):
        assert is_broken_vector([0.1, float("nan"), 0.2]) is True

    def test_normal_vector_is_fine(self):
        assert is_broken_vector([0.0, 0.3, -0.7, 1e-5]) is False


class TestStructuralChunker:
    def test_chunk_plain_text(self):
        extractor = ExtractedText(
            text="This is a test. " * 100,  # ~230 words
            pages=[],
            source_format="txt",
        )
        chunker = StructuralChunker(max_tokens=50)
        chunks = list(chunker.chunk(extractor))
        assert len(chunks) > 1
        assert all(c.chunk_type == "paragraph" for c in chunks)

    def test_chunk_preserves_tables(self):
        text = (
            "Some intro paragraph here.\n\n"
            "col1 col2 col3\n"
            "val1 val2 val3\n"
            "val4 val5 val6\n\n"
            "Another paragraph after the table."
        )
        extractor = ExtractedText(text=text, pages=[], source_format="txt")
        chunker = StructuralChunker(max_tokens=500)
        chunks = list(chunker.chunk(extractor))

        table_chunks = [c for c in chunks if c.chunk_type == "table"]
        assert len(table_chunks) >= 1
        assert "col1" in table_chunks[0].content

    def test_chunk_empty_text(self):
        extractor = ExtractedText(text="", pages=[], source_format="txt")
        chunker = StructuralChunker()
        chunks = list(chunker.chunk(extractor))
        assert len(chunks) == 0


class TestMetadataExtraction:
    def test_extract_title_from_content(self):
        text = "Housing Finance Guidelines\n\nSection 1: Credit scores..."
        extractor = ExtractedText(text=text, pages=[], source_format="txt")
        meta = extract_metadata(extractor, "doc.txt")
        assert meta.title is not None
        assert meta.doc_type == "policy"

    def test_extract_doc_type_from_content(self):
        text = "Underwriting guidelines for conventional loans..."
        extractor = ExtractedText(text=text, pages=[], source_format="txt")
        meta = extract_metadata(extractor, "underwriting.txt")
        assert meta.doc_type == "underwriting"


class TestEntityExtraction:
    def test_extracts_lenders(self):
        entities = extract_entities(
            "The FHA loan guidelines and VA loan program details."
        )
        assert "federal housing administration" in entities.lenders
        assert "veterans affairs" in entities.lenders

    def test_extracts_products(self):
        entities = extract_entities("An adjustable rate mortgage ARM or HELOC option.")
        assert "adjustable rate mortgage" in entities.products
        assert "helic home equity line of credit" in entities.products

    def test_extracts_documents(self):
        entities = extract_entities("Submit two pay stubs and a W-2 with your application.")
        assert "pay stub" in entities.documents
        assert "w-2" in entities.documents

    def test_extracts_property_types(self):
        entities = extract_entities("Financing a primary residence or an investment property.")
        assert "primary residence" in entities.property_types
        assert "investment property" in entities.property_types

    def test_empty_text_no_entities(self):
        entities = extract_entities("")
        assert entities.lenders == []
        assert entities.products == []
        assert entities.documents == []
        assert entities.property_types == []
        assert entities.acronyms == []


class TestTextExtractionFormats:
    def test_extract_csv(self, tmp_path):
        p = tmp_path / "data.csv"
        p.write_text("lender,limit\ngeneral,5000\nunderwriting,25000\n", encoding="utf-8")
        result = extract_text(p)
        assert "lender" in result.text
        assert "5000" in result.text
        assert "---" in result.text

    def test_extract_xlsx(self, tmp_path):
        import openpyxl

        p = tmp_path / "limits.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(["department", "approval_limit"])
        ws.append(["general", "5000"])
        wb.save(str(p))
        result = extract_text(p)
        assert "department" in result.text
        assert "5000" in result.text

    def test_extract_pptx(self, tmp_path):
        from pptx import Presentation

        p = tmp_path / "slides.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Approval Limits"
        prs.save(str(p))
        result = extract_text(p)
        assert "Approval Limits" in result.text

    def test_extract_rtf(self, tmp_path):
        p = tmp_path / "notes.rtf"
        p.write_bytes(b"{\\rtf1\\ansi Approval limit is \\b GBP 5000\\b0 .}")
        result = extract_text(p)
        assert "Approval limit" in result.text

    def test_extract_odt(self, tmp_path):
        import zipfile

        p = tmp_path / "doc.odt"
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<office:document-content '
            'xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" '
            'xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0" '
            'xmlns:table="urn:oasis:names:tc:opendocument:xmlns:table:1.0">'
            '<office:body><office:text>'
            "<text:p>Claims approval limit</text:p>"
            "<text:p>Escalate above limit</text:p>"
            "</office:text></office:body></office:document-content>"
        )
        with zipfile.ZipFile(str(p), "w") as zf:
            zf.writestr("content.xml", content_xml)
        result = extract_text(p)
        assert "Claims approval limit" in result.text

    def test_extract_epub(self, tmp_path):
        import zipfile

        p = tmp_path / "book.epub"
        chapter = '<html><body><p>Life insurance basics</p></body></html>'
        with zipfile.ZipFile(str(p), "w") as zf:
            zf.writestr("EPUB/chapter1.xhtml", chapter)
        result = extract_text(p)
        assert "Life insurance basics" in result.text

    def test_extract_htm(self, tmp_path):
        p = tmp_path / "page.htm"
        p.write_text("<html><body><h1>Guide</h1><p>Approval workflow.</p></body></html>", encoding="utf-8")
        result = extract_text(p)
        assert "Approval workflow" in result.text

    def test_unsupported_extension_raises(self, tmp_path):
        p = tmp_path / "archive.zip"
        p.write_bytes(b"not a document")
        try:
            extract_text(p)
        except ValueError:
            return
        raise AssertionError("expected ValueError for unsupported extension")


class TestZipBombCaps:
    """Audit 2.7: ODT/EPUB archive entries must be size-capped so a zip-bomb
    cannot inflate into the ~200MB-capped batch worker."""

    def test_odt_content_xml_over_entry_cap_rejected(self, tmp_path):
        import zipfile

        from unittest.mock import patch

        p = tmp_path / "bomb.odt"
        with zipfile.ZipFile(str(p), "w") as zf:
            zf.writestr("content.xml", "<text:p>" + "x" * 500 + "</text:p>")
        with patch("app.documents.text_extraction._MAX_ARCHIVE_ENTRY_BYTES", 100):
            with pytest.raises(ValueError, match="zip-bomb"):
                extract_text(p)

    def test_epub_chapter_over_entry_cap_rejected(self, tmp_path):
        import zipfile

        from unittest.mock import patch

        p = tmp_path / "bomb.epub"
        with zipfile.ZipFile(str(p), "w") as zf:
            zf.writestr("EPUB/chapter1.xhtml", "<p>" + "y" * 500 + "</p>")
        with patch("app.documents.text_extraction._MAX_ARCHIVE_ENTRY_BYTES", 100):
            with pytest.raises(ValueError, match="zip-bomb"):
                extract_text(p)

    def test_epub_total_decompressed_over_cap_rejected(self, tmp_path):
        import zipfile

        from unittest.mock import patch

        p = tmp_path / "manybomb.epub"
        with zipfile.ZipFile(str(p), "w") as zf:
            for i in range(4):
                zf.writestr(f"EPUB/c{i}.xhtml", "<p>" + "z" * 120 + "</p>")
        # 4 entries * ~126 bytes each < per-entry cap, but total > 200 cap.
        with (
            patch("app.documents.text_extraction._MAX_ARCHIVE_ENTRY_BYTES", 1000),
            patch("app.documents.text_extraction._MAX_ARCHIVE_TOTAL_BYTES", 200),
        ):
            with pytest.raises(ValueError, match="total cap"):
                extract_text(p)

    def test_legit_odt_and_epub_still_extract(self, tmp_path):
        import zipfile

        p = tmp_path / "ok.odt"
        content_xml = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<office:document-content '
            'xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" '
            'xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            '<office:body><office:text>'
            "<text:p>Loan to value ceiling</text:p>"
            "</office:text></office:body></office:document-content>"
        )
        with zipfile.ZipFile(str(p), "w") as zf:
            zf.writestr("content.xml", content_xml)
        assert "Loan to value ceiling" in extract_text(p).text

        p2 = tmp_path / "ok.epub"
        with zipfile.ZipFile(str(p2), "w") as zf:
            zf.writestr("EPUB/c1.xhtml", "<p>Term assurance basics</p>")
            zf.writestr("EPUB/c2.xhtml", "<p>Whole of life</p>")
        result = extract_text(p2)
        assert "Term assurance basics" in result.text
        assert "Whole of life" in result.text


class TestAllowedExtensions:
    def test_new_formats_accepted(self):
        for name in [
            "a.csv", "a.xlsx", "a.pptx", "a.odt", "a.rtf", "a.epub", "a.htm",
        ]:
            assert validate_upload(name, 1024).valid is True

    def test_unknown_extension_rejected(self):
        result = validate_upload("a.zip", 1024)
        assert result.valid is False
        assert "not allowed" in result.error
